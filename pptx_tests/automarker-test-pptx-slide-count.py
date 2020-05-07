from flask import Flask, request
from pptx import Presentation
from zappa.asynchronous import task

import common
import os
import requests
import sys

# Configuration and declaration
app = Flask(__name__)
# LAMBDA_API_KEY = 'testing'
###################################################
# You should not need to edit any of the code above


# Runs tests
# DEVELOPER TO-DO: This is where you will write your test
# Feel free to override the current test...
def run_tests(fileLocation, slideCountTestInfo):
    # DEVELOPER NOTES: This testsFailed array determines if a file passes a test
    # If it is returned as empty then a file is deemed to have passed the test
    # Else, the file fails to meet the criteria of the test...
    testsFailed = []
    path = str(fileLocation)

    #ensure cwd is in /tmp chdir
    os.chdir('/tmp')

    if (path.endswith('.pptx') == False):
        return ['FAIL: wrong file extension']
    
    try:
        prs = Presentation(path)
        if len(prs.slides) == 0:
            testsFailed = ['FAIL: There are no slides']
        else:
            print("Starting Test")
            print(slideCountTestInfo)
            slideCountTarget = 0 if 'target' not in slideCountTestInfo else slideCountTestInfo['target']
            slideCountMin = 1 if 'min' not in slideCountTestInfo else slideCountTestInfo['min']
            slideCountMax = 2 if 'max' not in slideCountTestInfo else slideCountTestInfo['max']
            slideCount = 0
            for slide in prs.slides:
                slideCount += 1
            if slideCountTarget > 0 and slideCount != slideCountTarget:
                testsFailed = ["FAIL: slideCount({}) didn't meet slideTarget({})".format(slideCount, slideCountTarget)]
            elif slideCountTarget <= 0 and slideCountMin > slideCount:
                testsFailed = ["FAIL: slideCount({}) is below slideCountMin({})".format(slideCount, slideCountMin)]
            elif slideCountTarget <= 0 and slideCountMax < slideCount:
                testsFailed = ["FAIL: slideCount({}) is above slideCountMax({})".format(slideCount, slideCountMax)]
            else:
                print("Passed. slideCount({})".format(slideCount))
            print(testsFailed)
    except Exception as e:
        print(e)
        testsFailed = ['FAIL: Something went wrong with parsing the pptx']

    return testsFailed

# Run async marking asynchronously
@task
def handle_async_marking(deliverable):
     # DEVELOPER TO-DO: This is where this endpoint extracts
    # meta data / parameters / arguments for the test
    # This test checks word count but a good word count can
    # be adjusted with a max and min count that we pass inside
    # the request to run this test
    # Feel free to override this part depending on the test you will write
    generalInfo = {} if 'generalAutoTestingInfo' not in deliverable else deliverable['generalAutoTestingInfo']
    slideCountTestInfo = {} if 'slideCountTest' not in generalInfo else generalInfo['slideCountTest']
    common.mark(deliverable, run_tests, slideCountTestInfo)

# You should not need to edit any code below
############################################
@app.route('/test', methods=['GET', 'POST'])
def begin_test():
   return common.begin_test(request, handle_async_marking)

if __name__ == '__main__':
    app.run(debug=True)
